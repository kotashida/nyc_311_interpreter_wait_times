
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import os

def create_charts():
    """Creates charts from the KPI data."""
    print("Creating charts...")
    if not os.path.exists("output/charts"):
        os.makedirs("output/charts")

    # Requests by language
    df = pd.read_csv("output/kpi_data/requests_by_language.csv")
    plt.figure()
    plt.bar(df['language'], df['request_count'])
    plt.title("Requests by Language")
    plt.savefig("output/charts/requests_by_language.png")

    # Median wait time
    df = pd.read_csv("output/kpi_data/median_wait_time.csv")
    plt.figure()
    plt.bar(df['language'], df['median_wait_time'])
    plt.title("Median Wait Time by Language")
    plt.savefig("output/charts/median_wait_time.png")

    # SLA compliance
    df = pd.read_csv("output/kpi_data/sla_compliance.csv")
    plt.figure()
    plt.bar(df['language'], df['sla_compliance_rate'])
    plt.title("SLA Compliance by Language")
    plt.savefig("output/charts/sla_compliance.png")

    # Borough breakdown
    df = pd.read_csv("output/kpi_data/borough_breakdown.csv")
    plt.figure()
    plt.bar(df['borough'], df['request_count'])
    plt.title("Requests by Borough")
    plt.xticks(rotation=45)
    plt.tight_layout()
    plt.savefig("output/charts/borough_breakdown.png")

    print("Charts created successfully.")

def create_presentation():
    """Creates a PowerPoint presentation from the charts."""
    print("Creating PowerPoint presentation...")
    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "NYC 311 Interpreter Access KPI Report"
    subtitle.text = "Weekly Summary"

    # Add charts
    for chart_file in os.listdir("output/charts"):
        if chart_file.endswith(".png"):
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = os.path.splitext(chart_file)[0].replace("_", " ").title()
            img_path = os.path.join("output/charts", chart_file)
            slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(8))

    prs.save("output/kpi_report.pptx")
    print("PowerPoint presentation created successfully.")

def run():
    """
    Generates a PowerPoint report from the KPI data.
    """
    create_charts()
    create_presentation()

if __name__ == "__main__":
    run()
